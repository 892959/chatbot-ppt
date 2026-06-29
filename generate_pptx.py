# -*- coding: utf-8 -*-
"""Generate chatbot分享 PPTX in NetEase customer-service template style."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
RED   = RGBColor(0xC0, 0x00, 0x00)
DARK  = RGBColor(0x1A, 0x1A, 0x1B)
INK   = RGBColor(0x40, 0x40, 0x40)
GREY  = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE= RGBColor(0xFF, 0x76, 0x00)
LRED  = RGBColor(0xF7, 0xE9, 0xE9)   # light red wash
LORG  = RGBColor(0xFF, 0xF1, 0xE6)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
LINEG = RGBColor(0xE5, 0xE5, 0xE5)
BLUE  = RGBColor(0x25, 0x63, 0xEB)

FONT = "Microsoft YaHei"
HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, "assets", "聚星工作台chatbot.png")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, round_=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.15, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
            ea = r.font._rPr.get_or_change_to_latin_typeface if False else None
    return tb

def footer(s, page):
    txt(s, Inches(0.55), Inches(7.02), Inches(7), Inches(0.4),
        [[("网易游戏 · 精灵运营部 · 客服业务", 10, GREY, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.4),
        [[(page, 10, GREY, False)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def header(s, num, cn, en, tag=None):
    # number badge
    if num:
        rect(s, Inches(0.55), Inches(0.55), Inches(0.62), Inches(0.62), fill=RED)
        txt(s, Inches(0.55), Inches(0.55), Inches(0.62), Inches(0.62),
            [[(num, 26, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = Inches(1.35)
    else:
        tx = Inches(0.55)
    txt(s, tx, Inches(0.48), Inches(11.3), Inches(0.62),
        [[(cn, 30, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tx, Inches(1.12), Inches(11.3), Inches(0.38),
        [[(en, 12.5, GREY, False)]])
    rect(s, tx, Inches(1.55), Inches(0.85), Inches(0.06), fill=RED)

def body_bg(s):
    rect(s, 0, 0, SW, SH, fill=WHITE)

# ============ 1. COVER ============
s = slide()
rect(s, 0, 0, SW, SH, fill=RED)
rect(s, 0, 0, SW, Inches(0.12), fill=DARK)
txt(s, Inches(0.9), Inches(0.6), Inches(8), Inches(0.4),
    [[("NETEASE GAMES · CUSTOMER SERVICE", 12, RGBColor(0xFF,0xCC,0xCC), False)]])
txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
    [[("Chatbot 相关的事项分享", 44, WHITE, True)]])
txt(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.5),
    [[("AI-Assisted Live Service Chatbot · Review & Methodology", 16, RGBColor(0xFF,0xDD,0xDD), False)]])
rect(s, Inches(0.95), Inches(4.45), Inches(1.1), Inches(0.07), fill=WHITE)
txt(s, Inches(0.9), Inches(4.7), Inches(11), Inches(0.5),
    [[("qzt    ·    精灵运营部月度分享    ·    2026.07", 15, RGBColor(0xFF,0xEE,0xEE), False)]])
txt(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.5),
    [[("—— 网易内部商业机密，请勿外传 ——", 12, RGBColor(0xFF,0xCC,0xCC), False)]])
txt(s, Inches(11.0), Inches(6.55), Inches(1.8), Inches(0.5),
    [[("網易 NETEASE", 13, WHITE, True)]], align=PP_ALIGN.RIGHT)

# ============ 2. CONTENTS ============
s = slide(); body_bg(s)
header(s, None, "目录", "Contents · Chatbot Sharing")
items = [("01","项目背景","机器人一环扣一环，chatbot 是最后一环"),
         ("02","困难与解法","采纳率为什么低？我们做了什么"),
         ("03","方法论沉淀","从 embedding 到「面粉与蛋糕」"),
         ("04","未来愿景","手动 → 半托管 → 全托管")]
y0 = Inches(2.1)
for i,(n,t,d) in enumerate(items):
    col = i % 2; row = i // 2
    x = Inches(0.9 + col*6.1); y = Emu(int(y0) + row*Inches(1.9))
    rect(s, x, y, Inches(0.06), Inches(1.55), fill=RED)
    txt(s, Emu(int(x)+Inches(0.25)), y, Inches(1.6), Inches(1.0),
        [[(n, 44, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Emu(int(x)+Inches(1.75)), Emu(int(y)+Inches(0.15)), Inches(4.2), Inches(0.5),
        [[(t, 21, INK, True)]])
    txt(s, Emu(int(x)+Inches(1.75)), Emu(int(y)+Inches(0.75)), Inches(4.2), Inches(0.7),
        [[(d, 12.5, GREY, False)]], line_spacing=1.2)
footer(s, "02")

# ============ 3. HOOK (quote) ============
s = slide(); rect(s, 0,0,SW,SH, fill=DARK)
txt(s, Inches(0.9), Inches(1.0), Inches(3), Inches(1.2),
    [[("\u201C", 80, RED, True)]])
txt(s, Inches(0.95), Inches(2.4), Inches(11.3), Inches(2.2),
    [[("一个玩家进线，背后是一条", 30, WHITE, True),("环环相扣", 30, ORANGE, True),("的 AI 流水线。", 30, WHITE, True)],
     [("而 chatbot，是直接面向客服的", 30, WHITE, True),("最后一环", 30, ORANGE, True),("。", 30, WHITE, True)]],
    line_spacing=1.4)
txt(s, Inches(0.95), Inches(4.9), Inches(11), Inches(0.6),
    [[("— 它能不能用、好不好用，决定了前面所有环节的价值能否兑现", 15, RGBColor(0xAA,0xAA,0xAA), False)]])
footer(s, "03")

# ============ 4. chatbot 长什么样 ============
s = slide(); body_bg(s)
header(s, "1", "Chatbot 长什么样", "In the agent workbench")
pts = [("1","嵌在聚星工作台会话流中","玩家话术下方实时浮现「AI 推荐回复」"),
       ("2","三种使用方式","一键发送 · 复制到编辑框 · 手动复制"),
       ("3","后验证拦截","不可用话术→拦截，显示「暂无话术推荐」"),
       ("4","0428 正式上线","覆盖 ma75(光遇) + xyq(梦幻西游) 等业务线")]
py = Inches(2.05)
for i,(n,t,d) in enumerate(pts):
    y = Emu(int(py) + i*Inches(1.12))
    c = rect(s, Inches(0.7), y, Inches(0.42), Inches(0.42), fill=RED, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.7), y, Inches(0.42), Inches(0.42), [[(n,13,WHITE,True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.3), Emu(int(y)-Inches(0.05)), Inches(5.0), Inches(0.45),
        [[(t,15.5,INK,True)]])
    txt(s, Inches(1.3), Emu(int(y)+Inches(0.42)), Inches(5.2), Inches(0.5),
        [[(d,11.5,GREY,False)]], line_spacing=1.15)
# right: AI card + image
rcx = Inches(6.9)
card = rect(s, rcx, Inches(1.95), Inches(5.7), Inches(1.35), fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(rcx)+Inches(0.2)), Inches(2.05), Inches(5), Inches(0.35),
    [[("🤖 AI 推荐回复", 11.5, RED, True)]])
txt(s, Emu(int(rcx)+Inches(0.2)), Inches(2.42), Inches(5.3), Inches(0.5),
    [[("\u201C理解大佬的心情，可以多留意后续活动哦~\u201D", 12, INK, False)]])
b1 = rect(s, Emu(int(rcx)+Inches(0.2)), Inches(2.85), Inches(2.0), Inches(0.36), fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(rcx)+Inches(0.2)), Inches(2.85), Inches(2.0), Inches(0.36), [[("复制并粘贴到输入框",10.5,WHITE,False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
b2 = rect(s, Emu(int(rcx)+Inches(2.35)), Inches(2.85), Inches(1.25), Inches(0.36), fill=CARD, line=RED, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(rcx)+Inches(2.35)), Inches(2.85), Inches(1.25), Inches(0.36), [[("一键发送",10.5,RED,False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
if os.path.exists(IMG):
    s.shapes.add_picture(IMG, rcx, Inches(3.45), width=Inches(5.7))
    txt(s, rcx, Inches(6.55), Inches(5.7), Inches(0.35),
        [[("聚星工作台真实界面 —— 红框处即 chatbot 推荐话术", 10.5, GREY, False)]], align=PP_ALIGN.CENTER)
footer(s, "04")

# ============ 5. 一环扣一环 ============
s = slide(); body_bg(s)
header(s, "1", "机器人一环扣一环", "An end-to-end AI service pipeline")
chain = [("VOC 意图识别","进人工前对话\n→ 三级意图", False),
         ("AX 工具执行","一个VOC调一个工具\n查知识库/数据/后台", False),
         ("Chatbot","工具结果+上下文\n→ 推荐话术", True),
         ("人工客服","一键发送 / 编辑 / 复制", False)]
nx = Inches(0.75); nw = Inches(2.7); gap = Inches(0.35); ny = Inches(2.3); nh = Inches(1.5)
for i,(t,d,hl) in enumerate(chain):
    x = Emu(int(nx) + i*(int(nw)+int(gap)))
    fill = LRED if hl else CARD
    ln = RED if hl else LINEG
    rect(s, x, ny, nw, nh, fill=fill, line=ln, line_w=2 if hl else 1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, Emu(int(ny)+Inches(0.22)), nw, Inches(0.45), [[(t,15,(RED if hl else INK),True)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(ny)+Inches(0.72)), nw, Inches(0.7), [[(d,11,GREY,False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 3:
        ax = Emu(int(x)+int(nw)+Inches(0.02))
        txt(s, ax, ny, gap, nh, [[("→",24,RED,True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
b = [("上游决定下游：","VOC 不准→工具调错→话术再漂亮也没用。chatbot 是最后一环，也是「兜底」与「放大器」。"),
     ("chatbot 的输入：","① AX 工具执行结果(主) ② 玩家与客服上下文 ③ 玩家当前发言 → 输出一句可直接用的推荐话术。")]
by = Inches(4.5)
for i,(h,d) in enumerate(b):
    y = Emu(int(by)+i*Inches(0.95))
    rect(s, Inches(0.75), Emu(int(y)+Inches(0.05)), Inches(0.32),Inches(0.32), fill=RED, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.75), Emu(int(y)+Inches(0.05)), Inches(0.32),Inches(0.32), [[("!",12,WHITE,True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.25), y, Inches(11.3), Inches(0.9),
        [[(h,13.5,INK,True),(d,13.5,GREY,False)]], line_spacing=1.2)
footer(s, "05")

# ============ 6. 上线成果 ============
s = slide(); body_bg(s)
header(s, "1", "上线成果速览", "Key metrics since 0428")
metrics = [("3.80%","辅助占比","理论上限 36.7%，约 9 倍空间", BLUE),
           ("66.61%","话术推荐率","有多少消息生成了推荐", ORANGE),
           ("55.10%","推荐可用率","推荐话术中可用的比例", RED)]
mx = Inches(0.9); mw = Inches(3.7); mg = Inches(0.35); my = Inches(2.3); mh = Inches(2.4)
for i,(num,lbl,sub,c) in enumerate(metrics):
    x = Emu(int(mx)+i*(int(mw)+int(mg)))
    rect(s, x, my, mw, mh, fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, Emu(int(my)+Inches(0.45)), mw, Inches(1.0), [[(num,46,c,True)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(my)+Inches(1.45)), mw, Inches(0.45), [[(lbl,17,INK,True)]], align=PP_ALIGN.CENTER)
    txt(s, x, Emu(int(my)+Inches(1.9)), mw, Inches(0.5), [[(sub,11.5,GREY,False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
txt(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.9),
    [[("从 5/8 的 2.03% 到 5/27 的 4.88%，闲聊能力、提示词优化每次迭代都带来肉眼可见的跃升 —— 但距离天花板，我们才走了不到 1/9。",14,GREY,False)]],
    align=PP_ALIGN.CENTER, line_spacing=1.3)
footer(s, "06")

# ============ 7. embedding 科普 ============
s = slide(); body_bg(s)
header(s, None, "大模型靠啥理解文字？", "Word Embedding · 词嵌入", )
txt(s, Inches(0.55), Inches(0.5), Inches(4), Inches(0.35), [[("📖 知识加油站",11,RED,True)]])
chips = [("充值","[0.82, -0.11, …]"),("没到账","[0.79, -0.09, …]"),("退款","[0.75, 0.20, …]"),("封号","[-0.30, 0.66, …]")]
cx=Inches(0.9)
for i,(w,v) in enumerate(chips):
    x=Emu(int(cx)+i*Inches(3.0))
    rect(s, x, Inches(1.95), Inches(2.7), Inches(0.8), fill=CARD, line=RED, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, Inches(2.02), Inches(2.7), Inches(0.4), [[(w,14,INK,True)]], align=PP_ALIGN.CENTER)
    txt(s, x, Inches(2.4), Inches(2.7), Inches(0.35), [[(v,11,RED,False)]], align=PP_ALIGN.CENTER)
b=[("文字 → 向量：","模型把每个词变成一串数字(向量)，意思相近的词向量也相近——“充值”和“没到账”挨得近，“封号”离得远。"),
   ("这是 VOC 识别的底层：","玩家说“我钱扣了东西没给我”，模型靠 embedding 判断它和“充值未到账”是一类，不需写死规则。"),
   ("这也是 chatbot 的底层：","理解上下文、匹配知识库、生成话术，全靠把语言压进同一个“语义空间”再计算距离。")]
by=Inches(3.15)
for i,(h,d) in enumerate(b):
    y=Emu(int(by)+i*Inches(0.9))
    rect(s, Inches(0.75), Emu(int(y)+Inches(0.04)),Inches(0.32),Inches(0.32),fill=RED,shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.75), Emu(int(y)+Inches(0.04)),Inches(0.32),Inches(0.32),[[(str(i+1),12,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.25), y, Inches(11.3),Inches(0.85),[[(h,13.5,INK,True),(d,13.5,GREY,False)]],line_spacing=1.2)
txt(s, Inches(0.9), Inches(6.3), Inches(11.5),Inches(0.5),
    [[("💡 一句话：embedding 把「语言」翻译成机器能算的「坐标」。",13,RED,True)]])
footer(s,"07")

# ============ 8. 三大瓶颈 ============
s = slide(); body_bg(s)
header(s, "2", "采纳率为什么上不去？", "Three bottlenecks")
cards=[("① 客服操作习惯","习惯自己打字，点选推荐的意愿待提升。工具再好，不用 = 0。"),
       ("② 上游 AX 结果","进人工只自动调一次工具，VOC 错了/玩家提新问题，结果就不匹配，话术失真。"),
       ("③ 话术拦截过严","后验证为保质量拦掉不少话术，可用推荐量被过滤过多。")]
cx=Inches(0.9); cw=Inches(3.7); cg=Inches(0.35); cy=Inches(2.3); ch=Inches(2.6)
for i,(h,d) in enumerate(cards):
    x=Emu(int(cx)+i*(int(cw)+int(cg)))
    rect(s, x, cy, cw, ch, fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, cy, cw, Inches(0.1), fill=RED)
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+Inches(0.35)), Emu(int(cw)-Inches(0.5)), Inches(0.6),[[(h,17,INK,True)]])
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+Inches(1.1)), Emu(int(cw)-Inches(0.5)), Inches(1.4),[[(d,13,GREY,False)]],line_spacing=1.3)
txt(s, Inches(0.9), Inches(5.4), Inches(11.5),Inches(0.6),
    [[("一个产品问题，",16,GREY,False),("一半在技术，一半在人",16,RED,True),("。",16,GREY,False)]],align=PP_ALIGN.CENTER)
footer(s,"08")

# ============ 9. 我们做了什么 ============
s = slide(); body_bg(s)
header(s, "2", "我们做了什么", "What we did & what's next")
cards=[("A 优化工作流","补充闲聊/安抚类话术；禁止“收到”“稍等”类废话。闲聊能力上线后曲线明显抬升。"),
       ("B 补高频知识库","针对高频场景补充识别与话术，提高推荐命中率。"),
       ("C 调拦截策略","保证质量前提下释放推荐量，平衡“可用率”与“覆盖率”。")]
cx=Inches(0.9); cw=Inches(3.7); cg=Inches(0.35); cy=Inches(2.2); ch=Inches(2.4)
for i,(h,d) in enumerate(cards):
    x=Emu(int(cx)+i*(int(cw)+int(cg)))
    rect(s, x, cy, cw, ch, fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, cy, cw, Inches(0.1), fill=ORANGE)
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+Inches(0.32)), Emu(int(cw)-Inches(0.5)), Inches(0.6),[[(h,17,INK,True)]])
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(cy)+Inches(1.05)), Emu(int(cw)-Inches(0.5)), Inches(1.3),[[(d,12.5,GREY,False)]],line_spacing=1.3)
rect(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.0), fill=LRED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.15), Inches(5.3), Inches(11.0), Inches(0.8),
    [[("→ 规划中：二次动态 VOC 总结调用 skill",14,RED,True)],
     [("不再只调一次工具——玩家意图变了，系统按「信息量」动态再识别一次、再调一次，自动纠偏。",12.5,GREY,False)]],line_spacing=1.25)
footer(s,"09")

# ============ 10. 面粉与蛋糕 quote ============
s = slide(); rect(s,0,0,SW,SH,fill=DARK)
txt(s, Inches(0.9), Inches(0.8), Inches(6), Inches(0.4),[[("03 · 方法论",12,ORANGE,True)]])
txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.8),
    [[("同样一袋面粉，",34,WHITE,True)],
     [("可以做成",34,WHITE,True),("馒头",34,ORANGE,True),("，也可以做成",34,WHITE,True),("蛋糕",34,ORANGE,True),("。",34,WHITE,True)]],line_spacing=1.35)
txt(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.0),
    [[("原料没变，关键看你怎么塑造、放在什么位置。",16,RGBColor(0xCC,0xCC,0xCC),False)],
     [("—— 同样的 AX 工具结果，提示词与方法不同，做出的话术天差地别；",16,RGBColor(0xCC,0xCC,0xCC),False)],
     [("—— 同样的客服，不该被钉死在“打字员”，他可以是盯盘的指挥官。",16,RGBColor(0xCC,0xCC,0xCC),False)]],line_spacing=1.4)
footer(s,"10")

# ============ 11. 比喻落地 ============
s = slide(); body_bg(s)
header(s,"3","原料相同，成果不同","Same flour, different cake")
# left card
lx=Inches(0.9); lw=Inches(5.4); ly=Inches(2.2); lh=Inches(3.4)
rect(s, lx, ly, lw, lh, fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(lx)+Inches(0.3)), Emu(int(ly)+Inches(0.25)), Emu(int(lw)-Inches(0.6)), Inches(0.5),[[("🌾 同样的“面粉”",18,INK,True)]])
for i,t in enumerate(["同一份 AX 工具结果","同一个大模型底座","同一批客服同事"]):
    txt(s, Emu(int(lx)+Inches(0.35)), Emu(int(ly)+Inches(1.0)+i*Inches(0.7)), Emu(int(lw)-Inches(0.7)), Inches(0.5),[[("=  "+t,14,GREY,False)]])
# right card
rx=Inches(6.95); rw=Inches(5.4)
rect(s, rx, ly, rw, lh, fill=LRED, line=RED, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(rx)+Inches(0.3)), Emu(int(ly)+Inches(0.25)), Emu(int(rw)-Inches(0.6)), Inches(0.5),[[("🎂 不同的“蛋糕”",18,RED,True)]])
for i,t in enumerate(["提示词/方法变通 → 话术可用率翻倍","流程变通 → 二次 VOC 自动纠偏","角色变通 → 客服从“回复”到“盯盘”"]):
    txt(s, Emu(int(rx)+Inches(0.35)), Emu(int(ly)+Inches(1.0)+i*Inches(0.7)), Emu(int(rw)-Inches(0.7)), Inches(0.5),[[("▲  "+t,13.5,INK,False)]])
txt(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.6),
    [[("基础资源是一样的，",14,GREY,False),("思路与方法不同，结果完全不同",14,RED,True),(" —— 这就是产品运营的价值。",14,GREY,False)]],align=PP_ALIGN.CENTER)
footer(s,"11")

# ============ 12. 方法论 ============
s = slide(); body_bg(s)
header(s,"3","沉淀下来的四条方法论","Methodology distilled")
b=[("先看上游，再看下游","采纳率低不一定是 chatbot 的锅——顺着 VOC→工具→话术追根因，别在最后一环死磕。"),
   ("用数据说话，小步快跑","每次优化(闲聊、提示词)都能在曲线上看到跳变。先建指标，再谈优化。"),
   ("技术问题，一半在人","客服不用，再好的推荐也是 0。降低操作成本、建立信任，和模型同等重要。"),
   ("学会变通","同样的原料，换个思路就是另一个成果。面粉可以是蛋糕，客服可以是指挥官。")]
by=Inches(2.15)
for i,(h,d) in enumerate(b):
    y=Emu(int(by)+i*Inches(1.12))
    rect(s, Inches(0.8), Emu(int(y)+Inches(0.05)),Inches(0.42),Inches(0.42),fill=RED,shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.8), Emu(int(y)+Inches(0.05)),Inches(0.42),Inches(0.42),[[(str(i+1),15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.4), Emu(int(y)-Inches(0.02)), Inches(11.0), Inches(0.45),[[(h,16,INK,True)]])
    txt(s, Inches(1.4), Emu(int(y)+Inches(0.45)), Inches(11.0), Inches(0.55),[[(d,12.5,GREY,False)]],line_spacing=1.2)
footer(s,"12")

# ============ 13. 未来愿景 ============
s = slide(); body_bg(s)
header(s,"4","从手动回复，到全托管","Manual → Co-pilot → Auto-pilot")
stages=[("现在","手动辅助","客服逐条点击「一键发送/复制」使用推荐", LRED, RED, INK),
        ("下一步","半托管","高置信话术自动发；低置信/有风险才提醒客服介入", LORG, ORANGE, INK),
        ("愿景","全托管 · 智能盯盘","机器人自动在发，没“堵车”不打扰；一个客服盯更多会话", DARK, ORANGE, WHITE)]
sx=Inches(0.9); sw=Inches(3.7); sg=Inches(0.35); sy=Inches(2.3); sh=Inches(2.6)
for i,(bd,t,d,bg,bc,tc) in enumerate(stages):
    x=Emu(int(sx)+i*(int(sw)+int(sg)))
    rect(s, x, sy, sw, sh, fill=bg, line=bc, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    badge=rect(s, Emu(int(x)+Inches(0.25)), Emu(int(sy)+Inches(0.25)), Inches(1.0), Inches(0.4), fill=bc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(sy)+Inches(0.25)), Inches(1.0), Inches(0.4),[[(bd,11,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, Emu(int(sy)+Inches(0.95)), sw, Inches(0.5),[[(t,17,tc,True)]],align=PP_ALIGN.CENTER)
    txt(s, Emu(int(x)+Inches(0.25)), Emu(int(sy)+Inches(1.55)), Emu(int(sw)-Inches(0.5)), Inches(0.9),[[(d,12,(RGBColor(0xCC,0xCC,0xCC) if i==2 else GREY),False)]],align=PP_ALIGN.CENTER,line_spacing=1.25)
    if i<2:
        txt(s, Emu(int(x)+int(sw)+Inches(0.02)), sy, sg, sh, [[("→",22,RED,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.7),
    [[("核心理念：没有卡点就自动流转，有“堵车”(准确率低/有风险/需抉择)才提醒人介入。",14,RED,True)]],align=PP_ALIGN.CENTER)
footer(s,"13")

# ============ 14. 好的问题比答案重要 ============
s = slide(); rect(s,0,0,SW,SH,fill=DARK)
txt(s, Inches(0.9), Inches(1.0), Inches(3), Inches(1.2),[[("\u201C",80,RED,True)]])
txt(s, Inches(0.95), Inches(2.5), Inches(11.3), Inches(1.0),
    [[("好的",36,WHITE,True),("问题",36,ORANGE,True),("，比答案更重要。",36,WHITE,True)]])
txt(s, Inches(0.95), Inches(3.8), Inches(11), Inches(0.5),[[("— 写在公司饭堂墙上的一句话",15,RGBColor(0xAA,0xAA,0xAA),False)]])
txt(s, Inches(0.95), Inches(4.7), Inches(11.3), Inches(1.5),
    [[("过去我们一直问：“怎么让话术更准？”但真正该问的也许是：",14,RGBColor(0xCC,0xCC,0xCC),False)],
     [("“客服到底需要什么样的辅助？”“我们用什么指标衡量它真的有用？”",14,RGBColor(0xCC,0xCC,0xCC),False)]],line_spacing=1.4)
footer(s,"14")

# ============ 15. 下半场 ============
s = slide(); body_bg(s)
header(s,None,"AI 的「下半场」","From how to solve to what's worth solving")
txt(s, Inches(0.55), Inches(0.5), Inches(5), Inches(0.35),[[("升华 · The Second Half",11,RED,True)]])
lx=Inches(0.9); lw=Inches(5.4); ly=Inches(2.2); lh=Inches(2.2)
rect(s, lx, ly, lw, lh, fill=CARD, line=LINEG, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(lx)+Inches(0.3)), Emu(int(ly)+Inches(0.3)), Emu(int(lw)-Inches(0.6)), Inches(0.5),[[("上半场",18,GREY,True)]])
txt(s, Emu(int(lx)+Inches(0.3)), Emu(int(ly)+Inches(1.0)), Emu(int(lw)-Inches(0.6)), Inches(1.0),[[("比拼怎么解题 —— 更强的模型、更高的准确率、更好的话术。",14,GREY,False)]],line_spacing=1.3)
rx=Inches(6.95); rw=Inches(5.4)
rect(s, rx, ly, rw, lh, fill=LRED, line=RED, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Emu(int(rx)+Inches(0.3)), Emu(int(ly)+Inches(0.3)), Emu(int(rw)-Inches(0.6)), Inches(0.5),[[("下半场",18,RED,True)]])
txt(s, Emu(int(rx)+Inches(0.3)), Emu(int(ly)+Inches(1.0)), Emu(int(rw)-Inches(0.6)), Inches(1.0),[[("比拼定义对的问题 —— 该解决什么、怎么衡量价值、如何重塑人与 AI 的协作。",14,INK,False)]],line_spacing=1.3)
txt(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.3),
    [[("姚顺雨说，AI 的下半场，瓶颈不再是“能不能做到”，而是“做什么、怎么定义成功”。",14,GREY,False)],
     [("chatbot 的下半场，不是更准的话术，而是 —— 重新定义客服与机器的关系。",15,RED,True)]],align=PP_ALIGN.CENTER,line_spacing=1.4)
footer(s,"15")

# ============ 16. 结尾 ============
s = slide(); rect(s,0,0,SW,SH,fill=RED)
rect(s, 0, 0, SW, Inches(0.12), fill=DARK)
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.8),
    [[("面粉一样，",40,WHITE,True)],[("但我们要做最好的那块蛋糕。",40,WHITE,True)]],line_spacing=1.3)
rect(s, Inches(0.95), Inches(4.5), Inches(1.1), Inches(0.07), fill=WHITE)
txt(s, Inches(0.9), Inches(4.75), Inches(11), Inches(0.5),[[("Thanks · Q&A    |    提问 · 讨论 · 一起把最后一环做扎实",16,RGBColor(0xFF,0xEE,0xEE),False)]])
txt(s, Inches(0.9), Inches(6.55), Inches(7), Inches(0.5),[[("网易游戏 · 精灵运营部 · 客服业务  |  邱梓涛",12,RGBColor(0xFF,0xDD,0xDD),False)]])
txt(s, Inches(0.9), Inches(6.95), Inches(11), Inches(0.4),[[("—— 网易内部商业机密，请勿外传 ——",11,RGBColor(0xFF,0xCC,0xCC),False)]])

out = os.path.join(HERE, "Chatbot相关的事项分享.pptx")
prs.save(out)
print("SAVED:", out)
